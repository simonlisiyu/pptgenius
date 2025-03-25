#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Author: lisiyu
Date: 2024
"""

import os
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import markdown
from bs4 import BeautifulSoup
from PIL import Image

class MarkdownToPPT:
    def __init__(self, template_path=None):
        """初始化转换器"""
        if template_path:
            # 如果提供了模板，先复制模板
            self.prs = Presentation(template_path)
            # 删除所有现有幻灯片
            for _ in range(len(self.prs.slides)):
                rId = self.prs.slides._sldIdLst[0].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[0]
        else:
            self.prs = Presentation()
        self.current_slide = None
        self.image_dir = None
        self.current_content_top = Inches(2)  # 从2英寸开始，给标题留空间

    def set_image_dir(self, image_dir):
        """设置图片目录"""
        self.image_dir = image_dir

    def add_slide(self):
        """添加新幻灯片"""
        # 创建空白幻灯片
        blank_layout = self.prs.slide_layouts[6]  # 使用空白布局
        self.current_slide = self.prs.slides.add_slide(blank_layout)
        self.current_content_top = Inches(2)  # 从2英寸开始，给标题留空间
        return self.current_slide

    def process_markdown(self, md_text):
        """处理Markdown内容"""
        # 配置Markdown解析器，启用表格和代码块扩展
        md = markdown.Markdown(extensions=['tables', 'fenced_code'])
        
        # 转换为HTML
        html = md.convert(md_text)
        
        # 使用BeautifulSoup解析HTML
        soup = BeautifulSoup(html, 'html.parser')
        
        # 处理每个元素
        current_slide = None
        for element in soup.find_all(['h1', 'h2', 'h3', 'p', 'img', 'table', 'pre', 'code', 'ul', 'li']):
            # 处理标题
            if element.name == 'h1':
                # 一级标题创建首页
                self.add_slide()
                self.add_title(element.get_text())
                self.current_content_top = Inches(2)  # 重置内容位置
            elif element.name == 'h2':
                # 二级标题创建新页
                self.add_slide()
                self.add_heading(element.get_text())
                self.current_content_top = Inches(1.5)  # 重置内容位置
            elif element.name == 'h3':
                # 三级标题作为子标题
                self.add_subheading(element.get_text())
            # 处理图片
            elif element.name == 'img':
                src = element.get('src')
                if src:
                    self.add_image(src)
            # 处理表格
            elif element.name == 'table':
                self.add_table(element)
            # 处理代码块
            elif element.name in ['pre', 'code']:
                self.add_code_block(element.get_text())
            # 处理列表项
            elif element.name == 'li':
                self.add_list_item(element.get_text())
            # 处理普通段落
            elif element.name == 'p' and not element.find_parent(['pre', 'code']):
                self.add_paragraph(element.get_text())

    def add_title(self, text):
        """添加首页标题"""
        if not self.current_slide:
            self.add_slide()
        
        # 计算标题位置和大小
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        
        # 首页标题居中显示
        left = Inches(0.5)
        top = slide_height / 3  # 垂直居中偏上
        width = slide_width - Inches(1)
        height = Inches(1.5)
        
        # 创建标题文本框
        title = self.current_slide.shapes.add_textbox(left, top, width, height)
        title.text_frame.text = text
        paragraph = title.text_frame.paragraphs[0]
        paragraph.font.size = Pt(44)  # 更大的字体
        paragraph.alignment = PP_ALIGN.CENTER

    def add_heading(self, text):
        """添加普通标题"""
        if not self.current_slide:
            self.add_slide()
        
        left = Inches(0.5)
        top = Inches(0.5)
        width = self.prs.slide_width - Inches(1)
        height = Inches(0.75)
        
        title = self.current_slide.shapes.add_textbox(left, top, width, height)
        title.text_frame.text = text
        paragraph = title.text_frame.paragraphs[0]
        paragraph.font.size = Pt(32)
        self.current_content_top = Inches(1.5)

    def add_subheading(self, text):
        """添加子标题"""
        if not self.current_slide:
            self.add_slide()
        
        left = Inches(0.5)
        width = self.prs.slide_width - Inches(1)
        height = Inches(0.5)
        
        title = self.current_slide.shapes.add_textbox(left, self.current_content_top, width, height)
        title.text_frame.text = text
        paragraph = title.text_frame.paragraphs[0]
        paragraph.font.size = Pt(24)
        self.current_content_top += Inches(0.75)

    def add_paragraph(self, text):
        """添加段落"""
        if not self.current_slide:
            self.add_slide()
        
        left = Inches(0.5)
        width = self.prs.slide_width - Inches(1)
        height = Inches(0.5)
        
        textbox = self.current_slide.shapes.add_textbox(left, self.current_content_top, width, height)
        textbox.text_frame.text = text
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.font.size = Pt(18)
        self.current_content_top += Inches(0.5)

    def add_list_item(self, text):
        """添加列表项"""
        if not self.current_slide:
            self.add_slide()
        
        left = Inches(1.0)  # 缩进
        width = self.prs.slide_width - Inches(1.5)
        height = Inches(0.5)
        
        textbox = self.current_slide.shapes.add_textbox(left, self.current_content_top, width, height)
        textbox.text_frame.text = f"• {text}"
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.font.size = Pt(18)
        self.current_content_top += Inches(0.5)

    def add_image(self, image_path):
        """添加图片"""
        if not self.current_slide:
            self.add_slide()
        
        print(f"\nDebug: Processing image: {image_path}")
        
        # 处理图片路径
        original_path = image_path
        if self.image_dir:
            image_path = os.path.join(self.image_dir, os.path.basename(image_path))
            print(f"Debug: Using image dir path: {image_path}")
        else:
            # 如果未指定图片目录，尝试在当前目录和images目录下查找
            possible_paths = [
                original_path,
                os.path.join('images', os.path.basename(original_path)),
                os.path.join(os.path.dirname(os.path.abspath(__file__)), 'images', os.path.basename(original_path))
            ]
            
            print("Debug: Trying possible paths:")
            for path in possible_paths:
                print(f"Debug: Checking path: {path}")
                if os.path.exists(path):
                    image_path = path
                    print(f"Debug: Found image at: {image_path}")
                    break
        
        if os.path.exists(image_path):
            try:
                print(f"Debug: Opening image file: {image_path}")
                # 获取幻灯片尺寸（转换为英寸）
                slide_width = Inches(10)  # 标准PPT宽度约为10英寸
                slide_height = Inches(7.5)  # 标准PPT高度约为7.5英寸
                
                # 计算图片尺寸
                img = Image.open(image_path)
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height
                print(f"Debug: Original image dimensions: {img_width} x {img_height}")
                
                # 设置图片最大尺寸（英寸）
                max_width = Inches(8)  # 最大宽度8英寸
                max_height = Inches(5)  # 最大高度5英寸
                
                # 计算最终尺寸（保持纵横比）
                if aspect_ratio > max_width/max_height:  # 宽图
                    width = max_width
                    height = width / aspect_ratio
                else:  # 高图
                    height = max_height
                    width = height * aspect_ratio
                
                print(f"Debug: Final dimensions in inches: {width/Inches(1)} x {height/Inches(1)}")
                
                # 计算居中位置
                left = (slide_width - width) / 2
                top = self.current_content_top
                print(f"Debug: Image position in inches: left={left/Inches(1)}, top={top/Inches(1)}")
                
                # 添加图片
                try:
                    shape = self.current_slide.shapes.add_picture(
                        image_path,
                        left, top,
                        width=width,
                        height=height
                    )
                    print(f"Debug: Successfully added image shape: {shape}")
                    # 更新下一个内容的位置
                    self.current_content_top += height + Inches(0.5)
                except Exception as e:
                    print(f"Error adding picture to slide: {str(e)}")
                    import traceback
                    print(traceback.format_exc())
            except Exception as e:
                print(f"Error processing image: {str(e)}")
                import traceback
                print(traceback.format_exc())
        else:
            print(f"Error: Image file not found: {image_path}")
            print(f"Current working directory: {os.getcwd()}")
            print(f"Image dir setting: {self.image_dir}")
            print(f"Absolute path attempted: {os.path.abspath(image_path)}")

    def add_table(self, table_element):
        """添加表格"""
        if not self.current_slide:
            self.add_slide()
        
        # 获取表格数据
        rows = []
        for tr in table_element.find_all('tr'):
            row = [td.get_text().strip() for td in tr.find_all(['td', 'th'])]
            rows.append(row)
        
        if not rows:
            return
        
        # 创建表格
        rows_count = len(rows)
        cols_count = len(rows[0])
        
        # 计算表格位置和大小
        left = Inches(0.5)
        width = self.prs.slide_width - Inches(1)
        height = Inches(0.5 * rows_count)
        
        # 添加表格
        table = self.current_slide.shapes.add_table(
            rows_count, cols_count,
            left, self.current_content_top,
            width, height
        ).table
        
        # 填充数据
        for i, row in enumerate(rows):
            for j, cell_text in enumerate(row):
                cell = table.cell(i, j)
                cell.text = cell_text
                cell.text_frame.paragraphs[0].font.size = Pt(14)
        
        self.current_content_top += height + Inches(0.5)

    def add_code_block(self, code):
        """添加代码块"""
        if not self.current_slide:
            self.add_slide()
        
        left = Inches(0.5)
        width = self.prs.slide_width - Inches(1)
        height = Inches(2)
        
        textbox = self.current_slide.shapes.add_textbox(left, self.current_content_top, width, height)
        textbox.text_frame.text = code
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.font.name = 'Courier New'
        paragraph.font.size = Pt(14)
        self.current_content_top += Inches(2.5)

    def save(self, output_path):
        """保存PPT文件"""
        self.prs.save(output_path)

def main():
    parser = argparse.ArgumentParser(description='Convert Markdown to PowerPoint')
    parser.add_argument('input_file', help='Input Markdown file')
    parser.add_argument('output_file', help='Output PowerPoint file')
    parser.add_argument('--template', help='PowerPoint template file')
    parser.add_argument('--image-dir', help='Directory containing images')
    
    args = parser.parse_args()
    
    # 读取Markdown文件
    with open(args.input_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    # 创建转换器实例
    converter = MarkdownToPPT(args.template)
    if args.image_dir:
        converter.set_image_dir(args.image_dir)
    
    # 处理Markdown内容
    converter.process_markdown(md_content)
    
    # 保存PPT
    converter.save(args.output_file)
    print(f'Successfully converted {args.input_file} to {args.output_file}')

if __name__ == '__main__':
    main() 